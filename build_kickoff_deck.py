from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Brand colors
NAVY = RGBColor(0x0B, 0x1D, 0x3A)
TEAL = RGBColor(0x00, 0x96, 0x9E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
ACCENT_GOLD = RGBColor(0xD4, 0xA0, 0x1E)
MEDIUM_GRAY = RGBColor(0x66, 0x66, 0x66)

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_bg(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multi_text(slide, left, top, width, height, lines, font_name="Calibri"):
    """lines = list of (text, font_size, color, bold, alignment)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, font_size, color, bold, alignment) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = alignment
        p.space_after = Pt(6)
    return txBox

def add_accent_bar(slide, left, top, width=Inches(1.5), height=Pt(4)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = TEAL
    shape.line.fill.background()
    return shape

# ==========================================
# SLIDE 1: Title Slide
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, NAVY)

# Accent bar at top
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Pt(6), TEAL)

# Bottom accent strip
add_shape_bg(slide, Inches(0), Inches(6.8), prs.slide_width, Inches(0.7), TEAL)

add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1),
    "SMARTWING™ TECHNOLOGY", 24, TEAL, True, PP_ALIGN.LEFT)

add_text_box(slide, Inches(1), Inches(2.3), Inches(11), Inches(1.5),
    "Internal Sales Kickoff 2026", 54, WHITE, True, PP_ALIGN.LEFT)

add_accent_bar(slide, Inches(1), Inches(4.0))

add_text_box(slide, Inches(1), Inches(4.4), Inches(11), Inches(1),
    "Tamarack Aerospace Group", 28, WHITE, False, PP_ALIGN.LEFT)

add_text_box(slide, Inches(1), Inches(6.9), Inches(5), Inches(0.5),
    "CONFIDENTIAL — INTERNAL USE ONLY", 14, NAVY, True, PP_ALIGN.LEFT)

add_text_box(slide, Inches(7), Inches(6.9), Inches(5.5), Inches(0.5),
    "220+ Installations Worldwide  |  FAA & EASA Certified", 14, NAVY, False, PP_ALIGN.RIGHT)

# ==========================================
# SLIDE 2: Agenda
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(4.5), prs.slide_height, NAVY)

add_text_box(slide, Inches(0.8), Inches(1), Inches(3), Inches(1),
    "AGENDA", 42, WHITE, True)
add_accent_bar(slide, Inches(0.8), Inches(1.9), Inches(1.2))

agenda_items = [
    ("01", "Where We Stand"),
    ("02", "SMARTWING™ Product Line"),
    ("03", "Model Performance Data"),
    ("04", "Sales Methodology"),
    ("05", "Buyer Personas"),
    ("06", "Objection Handling"),
    ("07", "New Tools & Resources"),
    ("08", "2026 Goals & Targets"),
]

for i, (num, title) in enumerate(agenda_items):
    y = Inches(2.5) + Inches(i * 0.55)
    add_text_box(slide, Inches(5.2), y, Inches(1), Inches(0.5),
        num, 22, TEAL, True)
    add_text_box(slide, Inches(6.0), y, Inches(6), Inches(0.5),
        title, 22, DARK_GRAY, False)

# ==========================================
# SLIDE 3: Where We Stand
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Pt(6), TEAL)

add_text_box(slide, Inches(1), Inches(0.5), Inches(11), Inches(0.8),
    "WHERE WE STAND", 38, WHITE, True)
add_accent_bar(slide, Inches(1), Inches(1.2))

# Stat boxes
stats = [
    ("220+", "Aircraft\nTransformed"),
    ("8", "Certified\nModels"),
    ("6", "Installation\nCenters"),
    ("5", "Countries\nApproved"),
]
for i, (number, label) in enumerate(stats):
    x = Inches(1) + Inches(i * 3.0)
    box = add_shape_bg(slide, x, Inches(2.0), Inches(2.6), Inches(2.2), RGBColor(0x10, 0x2A, 0x4F))
    add_text_box(slide, x, Inches(2.2), Inches(2.6), Inches(1),
        number, 52, TEAL, True, PP_ALIGN.CENTER)
    add_text_box(slide, x, Inches(3.2), Inches(2.6), Inches(0.8),
        label, 18, WHITE, False, PP_ALIGN.CENTER)

# Bottom text
add_multi_text(slide, Inches(1), Inches(4.8), Inches(11), Inches(2), [
    ("Key Milestones", 22, TEAL, True, PP_ALIGN.LEFT),
    ("• ForeFlight performance profile integration — fleet-wide", 16, WHITE, False, PP_ALIGN.LEFT),
    ("• ActiveCare Warranty program launched", 16, WHITE, False, PP_ALIGN.LEFT),
    ("• 200th installation milestone (Drew Forhan, CJ2+)", 16, WHITE, False, PP_ALIGN.LEFT),
    ("• Brazil expansion — first install with Solojet Aviação in São Paulo", 16, WHITE, False, PP_ALIGN.LEFT),
    ("• DaVinci Jets and Toledo Jet added to dealer network", 16, WHITE, False, PP_ALIGN.LEFT),
])

# ==========================================
# SLIDE 4: Product Line Overview
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.4), NAVY)

add_text_box(slide, Inches(1), Inches(0.3), Inches(11), Inches(0.8),
    "SMARTWING™ PRODUCT LINE", 36, WHITE, True)

products = [
    ("SMARTWING™", "Citation CJ / CJ1 / CJ1+ / CJ2 / CJ2+\nCJ3 / CJ3+ / M2 + GEN2 Variants", "Business Aviation"),
    ("Performance\nSMARTWING™", "King Air 200 / King Air 350", "Military & Special Mission"),
    ("Eco-\nSMARTWING™", "Commercial Aircraft Platforms", "Sustainability & Airlines"),
]
for i, (name, models, market) in enumerate(products):
    x = Inches(0.8) + Inches(i * 4.1)
    box = add_shape_bg(slide, x, Inches(2.0), Inches(3.7), Inches(4.5), LIGHT_GRAY)
    # Header bar
    add_shape_bg(slide, x, Inches(2.0), Inches(3.7), Inches(1.0), TEAL)
    add_text_box(slide, x, Inches(2.1), Inches(3.7), Inches(0.9),
        name, 22, WHITE, True, PP_ALIGN.CENTER)
    # Models
    add_text_box(slide, Inches(0.1) + x, Inches(3.3), Inches(3.5), Inches(1.5),
        models, 16, DARK_GRAY, False, PP_ALIGN.CENTER)
    # Market
    add_shape_bg(slide, x, Inches(5.2), Inches(3.7), Inches(0.6), NAVY)
    add_text_box(slide, x, Inches(5.25), Inches(3.7), Inches(0.5),
        market, 14, WHITE, True, PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(6.8), Inches(11), Inches(0.5),
    "All products built on patented ATLAS® load alleviation technology  •  FAA & EASA Certified", 14, MEDIUM_GRAY, False, PP_ALIGN.CENTER)

# ==========================================
# SLIDE 5: Performance Data — CJ/CJ1/CJ1+
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.4), NAVY)

add_text_box(slide, Inches(1), Inches(0.3), Inches(11), Inches(0.8),
    "PERFORMANCE DATA — CJ / CJ1 / CJ1+", 36, WHITE, True)

# Table header
headers = ["Metric", "Stock CJ", "SMARTWING™ CJ", "Delta"]
col_x = [Inches(1), Inches(4), Inches(7), Inches(10.5)]
col_w = [Inches(3), Inches(3), Inches(3.5), Inches(2)]

add_shape_bg(slide, Inches(0.8), Inches(1.8), Inches(11.7), Inches(0.6), TEAL)
for j, h in enumerate(headers):
    add_text_box(slide, col_x[j], Inches(1.85), col_w[j], Inches(0.5),
        h, 16, WHITE, True, PP_ALIGN.CENTER)

rows = [
    ("MZFW", "8,400 lb", "8,800 lb", "+400 lb"),
    ("Max Payload", "1,730 lb", "2,052 lb", "+322 lb"),
    ("Time to FL410", "60+ min", "≤32 min", "-28 min"),
    ("Block Fuel Burn", "900 pph", "600–650 pph", "-28%"),
    ("Range @ MCT", "1,100 nm / 3 hr", "1,400 nm / 4 hr", "+300 nm"),
]
for i, row in enumerate(rows):
    y = Inches(2.5) + Inches(i * 0.55)
    bg_color = LIGHT_GRAY if i % 2 == 0 else WHITE
    add_shape_bg(slide, Inches(0.8), y, Inches(11.7), Inches(0.55), bg_color)
    for j, cell in enumerate(row):
        c = TEAL if j == 3 else DARK_GRAY
        b = True if j == 3 else False
        add_text_box(slide, col_x[j], y + Pt(4), col_w[j], Inches(0.45),
            cell, 15, c, b, PP_ALIGN.CENTER)

# Key callout
add_shape_bg(slide, Inches(1), Inches(5.5), Inches(11), Inches(1.5), NAVY)
add_multi_text(slide, Inches(1.5), Inches(5.6), Inches(10), Inches(1.4), [
    ("CJ WORLD RECORD: 1,853 nm nonstop (6 hr 16 min with 26 kt tailwind)", 20, ACCENT_GOLD, True, PP_ALIGN.LEFT),
    ("High-hot: +549 lbs additional takeoff weight at 5,000 ft / 44°C", 16, WHITE, False, PP_ALIGN.LEFT),
    ("OEI climb gradient weights up to 1,000 lb higher  •  Airport noise reduction certified", 16, WHITE, False, PP_ALIGN.LEFT),
])

# ==========================================
# SLIDE 6: Performance Data — CJ2/CJ2+
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.4), NAVY)

add_text_box(slide, Inches(1), Inches(0.3), Inches(11), Inches(0.8),
    "PERFORMANCE DATA — CJ2 / CJ2+", 36, WHITE, True)

add_shape_bg(slide, Inches(0.8), Inches(1.8), Inches(11.7), Inches(0.6), TEAL)
for j, h in enumerate(headers):
    add_text_box(slide, col_x[j], Inches(1.85), col_w[j], Inches(0.5),
        h, 16, WHITE, True, PP_ALIGN.CENTER)

rows = [
    ("MZFW", "9,300 lb", "10,100 lb", "+800 lb"),
    ("Max Payload", "1,490 lb", "2,211 lb", "+721 lb"),
    ("Time to FL450", "40+ min", "≤30 min", "-10 min"),
    ("Block Fuel Burn", "900 pph", "680–700 pph", "-22%"),
    ("Range @ MCT", "1,200 nm / 4.5 hr", "1,650 nm / 5.5 hr", "+450 nm"),
]
for i, row in enumerate(rows):
    y = Inches(2.5) + Inches(i * 0.55)
    bg_color = LIGHT_GRAY if i % 2 == 0 else WHITE
    add_shape_bg(slide, Inches(0.8), y, Inches(11.7), Inches(0.55), bg_color)
    for j, cell in enumerate(row):
        c = TEAL if j == 3 else DARK_GRAY
        b = True if j == 3 else False
        add_text_box(slide, col_x[j], y + Pt(4), col_w[j], Inches(0.45),
            cell, 15, c, b, PP_ALIGN.CENTER)

add_shape_bg(slide, Inches(1), Inches(5.5), Inches(11), Inches(1.5), NAVY)
add_multi_text(slide, Inches(1.5), Inches(5.6), Inches(10), Inches(1.4), [
    ("THE ALTITUDE ADVANTAGE", 20, ACCENT_GOLD, True, PP_ALIGN.LEFT),
    ("Stock CJ2/CJ2+ at FL400 operates at 16% lower efficiency than SMARTWING™ at FL450", 16, WHITE, False, PP_ALIGN.LEFT),
    ("CJ2+ payload jumps from 1,700 lb to 2,425 lb (+725 lb)  •  MZFW +400 lb", 16, WHITE, False, PP_ALIGN.LEFT),
    ("Drew Forhan — 200th installation, CJ2+ owner, Tamarack's most visible advocate", 16, WHITE, False, PP_ALIGN.LEFT),
])

# ==========================================
# SLIDE 7: Performance Data — CJ3 & M2
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.4), NAVY)

add_text_box(slide, Inches(1), Inches(0.3), Inches(11), Inches(0.8),
    "PERFORMANCE DATA — CJ3 / CJ3+ & M2", 36, WHITE, True)

# CJ3 mini-table
add_text_box(slide, Inches(1), Inches(1.7), Inches(5), Inches(0.5),
    "CJ3 / CJ3+", 20, TEAL, True)
cj3_rows = [
    ("Range", "+250 nm (1,900 nm)"),
    ("Fuel Burn", "-22% (680 pph)"),
    ("Payload", "+320 lb (2,645 lb)"),
    ("Climb", "FL450 in ≤25 min"),
]
for i, (metric, value) in enumerate(cj3_rows):
    y = Inches(2.3) + Inches(i * 0.45)
    bg = LIGHT_GRAY if i % 2 == 0 else WHITE
    add_shape_bg(slide, Inches(1), y, Inches(5.5), Inches(0.45), bg)
    add_text_box(slide, Inches(1.2), y + Pt(2), Inches(2), Inches(0.4),
        metric, 15, DARK_GRAY, True, PP_ALIGN.LEFT)
    add_text_box(slide, Inches(3.2), y + Pt(2), Inches(3), Inches(0.4),
        value, 15, TEAL, True, PP_ALIGN.LEFT)

# M2 mini-table
add_text_box(slide, Inches(7.2), Inches(1.7), Inches(5), Inches(0.5),
    "Citation M2", 20, TEAL, True)
m2_rows = [
    ("Range", "+150 nm (1,700 nm)"),
    ("Fuel Burn", "-13% (690 pph)"),
    ("Payload", "+323 lb (1,823 lb)"),
    ("Climb", "FL410 in ≤22 min"),
]
for i, (metric, value) in enumerate(m2_rows):
    y = Inches(2.3) + Inches(i * 0.45)
    bg = LIGHT_GRAY if i % 2 == 0 else WHITE
    add_shape_bg(slide, Inches(7.2), y, Inches(5.5), Inches(0.45), bg)
    add_text_box(slide, Inches(7.4), y + Pt(2), Inches(2), Inches(0.4),
        metric, 15, DARK_GRAY, True, PP_ALIGN.LEFT)
    add_text_box(slide, Inches(9.4), y + Pt(2), Inches(3), Inches(0.4),
        value, 15, TEAL, True, PP_ALIGN.LEFT)

# CJ3 callout
add_shape_bg(slide, Inches(1), Inches(4.5), Inches(5.5), Inches(2.5), NAVY)
add_multi_text(slide, Inches(1.3), Inches(4.6), Inches(5), Inches(2.3), [
    ("CJ3 FLIGHT STORIES", 18, ACCENT_GOLD, True, PP_ALIGN.LEFT),
    ("• Shannon → Gander transatlantic crossing", 14, WHITE, False, PP_ALIGN.LEFT),
    ("• Mosier world trip", 14, WHITE, False, PP_ALIGN.LEFT),
    ("• Orlando round-trip case study", 14, WHITE, False, PP_ALIGN.LEFT),
    ("• Transatlantic capability is a CJ3-unique selling point", 14, WHITE, False, PP_ALIGN.LEFT),
])

# M2 callout
add_shape_bg(slide, Inches(7.2), Inches(4.5), Inches(5.5), Inches(2.5), NAVY)
add_multi_text(slide, Inches(7.5), Inches(4.6), Inches(5), Inches(2.3), [
    ("M2 EXCLUSIVE", 18, ACCENT_GOLD, True, PP_ALIGN.LEFT),
    ("• YAW Damper INOP restriction REMOVED", 14, WHITE, False, PP_ALIGN.LEFT),
    ("• Unique benefit — no other model has this", 14, WHITE, False, PP_ALIGN.LEFT),
    ("• GEN2 variant now certified", 14, WHITE, False, PP_ALIGN.LEFT),
    ("• Strongest aesthetic transformation in the lineup", 14, WHITE, False, PP_ALIGN.LEFT),
])

# ==========================================
# SLIDE 8: SPIN Selling Methodology
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Pt(6), TEAL)

add_text_box(slide, Inches(1), Inches(0.5), Inches(11), Inches(0.8),
    "SALES METHODOLOGY — SPIN FRAMEWORK", 36, WHITE, True)
add_accent_bar(slide, Inches(1), Inches(1.2))

spin = [
    ("S", "SITUATION", "Establish the prospect's aircraft, usage, and mission profile",
     '"What model? How many hours/year? Typical mission distance?"'),
    ("P", "PROBLEM", "Surface pain points they may not have articulated",
     '"Ever had to make a fuel stop you didn\'t want to? Limited at high-altitude airports?"'),
    ("I", "IMPLICATION", "Deepen the pain — make the cost of inaction tangible",
     '"What does that extra fuel stop cost in time, fuel, and wear?"'),
    ("N", "NEED-PAYOFF", "Let the prospect articulate the value of the solution",
     '"If you could eliminate that fuel stop, what would it mean for your schedule?"'),
]

for i, (letter, title, desc, example) in enumerate(spin):
    y = Inches(1.8) + Inches(i * 1.3)
    # Letter circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1), y, Inches(0.8), Inches(0.8))
    circle.fill.solid()
    circle.fill.fore_color.rgb = TEAL
    circle.line.fill.background()
    add_text_box(slide, Inches(1), y + Pt(6), Inches(0.8), Inches(0.6),
        letter, 32, WHITE, True, PP_ALIGN.CENTER)
    
    add_text_box(slide, Inches(2.1), y, Inches(3), Inches(0.4),
        title, 22, TEAL, True)
    add_text_box(slide, Inches(2.1), y + Inches(0.35), Inches(4.5), Inches(0.5),
        desc, 14, WHITE, False)
    add_text_box(slide, Inches(7), y + Inches(0.05), Inches(5.5), Inches(0.9),
        example, 13, RGBColor(0xAA, 0xAA, 0xAA), False, PP_ALIGN.LEFT)

# ==========================================
# SLIDE 9: Buyer Personas
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.4), NAVY)

add_text_box(slide, Inches(1), Inches(0.3), Inches(11), Inches(0.8),
    "BUYER PERSONAS", 36, WHITE, True)

personas = [
    ("THE GO-GETTER", "Self-made, 45–65\nTime is money", "Lead with ROI\nand performance data", TEAL),
    ("THE COLLECTOR", "Affluent, multiple aircraft\nValues craftsmanship", "Lead with Vref,\naesthetics, pedigree", RGBColor(0x2E, 0x86, 0xAB)),
    ("THE OLD PRO", "Career aviator, 55+\nThousands of hours", "Go deep on aero.\nLet data speak.", RGBColor(0x1B, 0x4F, 0x72)),
    ("THE SKEPTIC", "Analytical, research-driven\nFinds every forum post", "Address AD proactively.\nProvide documentation.", RGBColor(0x0B, 0x2D, 0x4A)),
]

for i, (name, demo, approach, color) in enumerate(personas):
    x = Inches(0.5) + Inches(i * 3.2)
    # Card
    add_shape_bg(slide, x, Inches(1.8), Inches(2.9), Inches(5.0), LIGHT_GRAY)
    # Header
    add_shape_bg(slide, x, Inches(1.8), Inches(2.9), Inches(1.0), color)
    add_text_box(slide, x, Inches(1.95), Inches(2.9), Inches(0.8),
        name, 18, WHITE, True, PP_ALIGN.CENTER)
    # Demo
    add_text_box(slide, Inches(0.15) + x, Inches(3.0), Inches(2.6), Inches(1.2),
        "DEMOGRAPHICS", 11, MEDIUM_GRAY, True, PP_ALIGN.LEFT)
    add_text_box(slide, Inches(0.15) + x, Inches(3.35), Inches(2.6), Inches(1.0),
        demo, 14, DARK_GRAY, False, PP_ALIGN.LEFT)
    # Approach
    add_text_box(slide, Inches(0.15) + x, Inches(4.5), Inches(2.6), Inches(0.4),
        "APPROACH", 11, MEDIUM_GRAY, True, PP_ALIGN.LEFT)
    add_text_box(slide, Inches(0.15) + x, Inches(4.85), Inches(2.6), Inches(1.5),
        approach, 14, DARK_GRAY, False, PP_ALIGN.LEFT)

# ==========================================
# SLIDE 10: Top Objections
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Pt(6), TEAL)

add_text_box(slide, Inches(1), Inches(0.5), Inches(11), Inches(0.8),
    "TOP 6 OBJECTIONS & RESPONSES", 36, WHITE, True)
add_accent_bar(slide, Inches(1), Inches(1.2))

objections = [
    ('"Too expensive"', "Vref 100% + fuel savings = payback in 1–2 years"),
    ('"I heard about the AD"', "Acknowledge → explain resolution → pivot to current fleet performance"),
    ('"I\'m selling the aircraft"', "SMARTWING™ increases resale value — install first, sell for more"),
    ('"My mechanic doesn\'t like it"', "Connect them directly with Tamarack engineering"),
    ('"I don\'t fly enough"', "Even at 100 hrs/yr, Vref alone justifies it. Safety gains aren't usage-dependent"),
    ('"Send me information"', "\"What's the ONE thing you'd want to see?\" Then schedule follow-up in 3 days"),
]

for i, (objection, response) in enumerate(objections):
    y = Inches(1.7) + Inches(i * 0.9)
    add_shape_bg(slide, Inches(1), y, Inches(4.2), Inches(0.7), RGBColor(0x10, 0x2A, 0x4F))
    add_text_box(slide, Inches(1.3), y + Pt(6), Inches(3.8), Inches(0.6),
        objection, 16, ACCENT_GOLD, True, PP_ALIGN.LEFT)
    add_text_box(slide, Inches(5.5), y + Pt(6), Inches(7), Inches(0.6),
        response, 15, WHITE, False, PP_ALIGN.LEFT)

add_text_box(slide, Inches(1), Inches(7.0), Inches(11), Inches(0.4),
    "Full 13-objection playbook + 8 long-term owner objections available in Sales Kit documents", 13, MEDIUM_GRAY, False, PP_ALIGN.CENTER)

# ==========================================
# SLIDE 11: New Tools & Resources
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.4), NAVY)

add_text_box(slide, Inches(1), Inches(0.3), Inches(11), Inches(0.8),
    "NEW TOOLS IN YOUR SALES KIT", 36, WHITE, True)

tools = [
    ("Model One-Pagers", "One-page leave-behinds for CJ, CJ1/CJ1+, CJ2,\nCJ2+, CJ3/CJ3+, and M2 — 3 key numbers,\ncustomer proof, CTA", "Sales and Marketing/One-Pagers/"),
    ("Follow-Up Sequence", "60-day structured follow-up cadence with\nemail templates, call scripts, and persona-\nmatched content at every touchpoint", "Sales and Marketing/Sales Tools/"),
    ("Competitive Positioning", "SMARTWING™ vs. 6 alternatives: do nothing,\ntrade up, passive winglets, other mods,\nfuture tech, buy pre-equipped", "Sales and Marketing/Sales Tools/"),
    ("ActiveCare & ForeFlight", "Talking points for ActiveCare Warranty and\nForeFlight performance profile integration —\ntwo selling tools you're not using yet", "Sales and Marketing/Sales Tools/"),
    ("Sustainability Sell Sheet", "ESG-ready pitch for corporate flight departments\nwith CO2 data, ICAO credentials, and\nregulatory context", "Sales and Marketing/One-Pagers/"),
    ("GEN2 Sales Guide", "M2 GEN2, CJ3 GEN2, CJ3+ GEN2 —\nprospecting strategy and objection handling\nfor newest-model buyers", "Aircraft Models/"),
]

for i, (title, desc, loc) in enumerate(tools):
    col = i % 2
    row = i // 2
    x = Inches(0.8) + Inches(col * 6.2)
    y = Inches(1.8) + Inches(row * 1.7)
    add_shape_bg(slide, x, y, Inches(5.8), Inches(1.5), LIGHT_GRAY)
    add_shape_bg(slide, x, y, Inches(0.08), Inches(1.5), TEAL)
    add_text_box(slide, x + Inches(0.3), y + Pt(4), Inches(5.2), Inches(0.4),
        title, 18, NAVY, True)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.45), Inches(5.2), Inches(0.8),
        desc, 12, DARK_GRAY, False)
    add_text_box(slide, x + Inches(0.3), y + Inches(1.15), Inches(5.2), Inches(0.3),
        loc, 10, MEDIUM_GRAY, False)

# ==========================================
# SLIDE 12: The ROI Story
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Pt(6), TEAL)

add_text_box(slide, Inches(1), Inches(0.5), Inches(11), Inches(0.8),
    "THE 3-PART ROI STORY", 36, WHITE, True)
add_accent_bar(slide, Inches(1), Inches(1.2))

roi_parts = [
    ("1", "VREF VALUE", "100% resale value at installation.\n90% at year 2.\n\n\"This is the only mod where you\nget your money back before\nyou even fly it.\""),
    ("2", "SOFT VALUE", "8 categories rated 1–5:\n\nBuying habits • Aesthetics\nSafety • Comfort\nDispatch • Environmental\nRange • Payload\n\nTypical: 20–30% above install price"),
    ("3", "FUEL SAVINGS", "Real-world calculation:\n\nCJ2 @ 300 hrs/yr:\nStock: 270,000 lbs fuel\nSMARTWING™: 207,000 lbs\nSavings: ~$51,800/year\n\nPayback: Under 4 years"),
]

for i, (num, title, body) in enumerate(roi_parts):
    x = Inches(0.8) + Inches(i * 4.1)
    add_shape_bg(slide, x, Inches(1.8), Inches(3.8), Inches(5.0), RGBColor(0x10, 0x2A, 0x4F))
    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(1.5), Inches(2.0), Inches(0.7), Inches(0.7))
    circle.fill.solid()
    circle.fill.fore_color.rgb = TEAL
    circle.line.fill.background()
    add_text_box(slide, x + Inches(1.5), Inches(2.05), Inches(0.7), Inches(0.6),
        num, 28, WHITE, True, PP_ALIGN.CENTER)
    
    add_text_box(slide, x + Inches(0.3), Inches(2.9), Inches(3.2), Inches(0.5),
        title, 20, ACCENT_GOLD, True, PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.3), Inches(3.5), Inches(3.2), Inches(3.0),
        body, 14, WHITE, False, PP_ALIGN.CENTER)

# ==========================================
# SLIDE 13: 2026 Priorities
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.4), NAVY)

add_text_box(slide, Inches(1), Inches(0.3), Inches(11), Inches(0.8),
    "2026 PRIORITIES", 36, WHITE, True)

priorities = [
    ("Fill the proof gap", "Collect FlightAware tracks and testimonials for CJ1+ and CJ2 — our two weakest models for proof assets"),
    ("Leverage ForeFlight + ActiveCare", "Every prospect call should mention ForeFlight integration and ActiveCare Warranty — these are differentiators we're not using"),
    ("Build the charter pipeline", "CJ2 and CJ3 charter operators have the fastest ROI and highest utilization — develop charter-specific materials and target Part 135 operators"),
    ("GEN2 early movers", "M2 GEN2, CJ3 GEN2, and CJ3+ GEN2 buyers are upgrade-minded — reach them during the purchase process"),
    ("Sustainability positioning", "Corporate flight departments with ESG mandates are a growing segment — use ICAO credentials and CO2 data"),
    ("Follow-up discipline", "Use the new 60-day follow-up sequence for every prospect. Log every touch in CRM. No more deals lost to silence."),
]

for i, (title, desc) in enumerate(priorities):
    col = i % 2
    row = i // 2
    x = Inches(0.8) + Inches(col * 6.2)
    y = Inches(1.8) + Inches(row * 1.7)
    
    # Number
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y + Inches(0.15), Inches(0.5), Inches(0.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = TEAL
    circle.line.fill.background()
    add_text_box(slide, x, y + Inches(0.18), Inches(0.5), Inches(0.45),
        str(i + 1), 20, WHITE, True, PP_ALIGN.CENTER)
    
    add_text_box(slide, x + Inches(0.7), y, Inches(5.2), Inches(0.4),
        title, 18, NAVY, True)
    add_text_box(slide, x + Inches(0.7), y + Inches(0.4), Inches(5.2), Inches(1.0),
        desc, 13, DARK_GRAY, False)

# ==========================================
# SLIDE 14: Closing Slide
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_shape_bg(slide, Inches(0), Inches(0), prs.slide_width, Pt(6), TEAL)
add_shape_bg(slide, Inches(0), Inches(6.8), prs.slide_width, Inches(0.7), TEAL)

add_text_box(slide, Inches(1), Inches(2.0), Inches(11.3), Inches(1.5),
    "Let's make 2026 our\nbiggest year yet.", 52, WHITE, True, PP_ALIGN.CENTER)

add_accent_bar(slide, Inches(5.9), Inches(4.0), Inches(1.5))

add_text_box(slide, Inches(2), Inches(4.5), Inches(9.3), Inches(0.6),
    "220+ installations  •  8 models certified  •  6 global centers", 22, TEAL, False, PP_ALIGN.CENTER)

add_text_box(slide, Inches(2), Inches(5.3), Inches(9.3), Inches(0.5),
    "tamarackaero.com", 20, WHITE, False, PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(6.9), Inches(11.3), Inches(0.5),
    "CONFIDENTIAL — INTERNAL USE ONLY  |  TAMARACK AEROSPACE GROUP  |  2026", 13, NAVY, True, PP_ALIGN.CENTER)

# Save
output_path = r"c:\Users\honeybadger\Dev\Repos\Sales\Presentations\SMARTWING Internal Sales Kickoff 2026.pptx"
prs.save(output_path)
print(f"Saved to {output_path}")
